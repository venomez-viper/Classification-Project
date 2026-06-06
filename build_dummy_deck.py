"""
Build a reference dummy deck in the Dark Minimalist palette.
Output: REFERENCE_DECK.pptx in the capstone folder.
Use this as a content source — copy text from here into the real template.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# Palette extracted from Dark Minimalist Presentation.pptx
BG     = RGBColor(0x00, 0x15, 0x14)  # #001514 dark teal-black
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LIME   = RGBColor(0xC2, 0xD0, 0x76)  # #C2D076 olive-lime accent
CORAL  = RGBColor(0xFF, 0x6B, 0x6B)  # #FF6B6B leakage warning
DIM    = RGBColor(0x6B, 0x70, 0x70)  # muted gray for sub-captions

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

BLANK = prs.slide_layouts[6]  # truly blank

def set_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text(slide, text, left, top, width, height, *,
             size=24, color=WHITE, bold=False, italic=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    # First paragraph already exists
    paragraphs = text.split("\n")
    for i, line in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.italic = italic
    return box

def add_bullets(slide, items, left, top, width, height, *,
                size=22, color=WHITE, bold_keys=None, font="Calibri"):
    """items: list of strings. Each rendered as a bullet line."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)
        # Custom bullet: lime square
        bullet = p.add_run()
        bullet.text = "■  "
        bullet.font.name = font
        bullet.font.size = Pt(size)
        bullet.font.color.rgb = LIME
        body = p.add_run()
        body.text = item
        body.font.name = font
        body.font.size = Pt(size)
        body.font.color.rgb = color
    return box

def add_rule(slide, left, top, width, color=LIME, thickness_pt=2):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(thickness_pt))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    return line

def add_footer(slide, text, lime_part=None):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.4))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = text
    r.font.name = "Calibri"
    r.font.size = Pt(11)
    r.font.color.rgb = DIM
    if lime_part:
        r2 = p.add_run()
        r2.text = "   " + lime_part
        r2.font.name = "Calibri"
        r2.font.size = Pt(11)
        r2.font.color.rgb = LIME

def add_notes(slide, notes_text):
    notes = slide.notes_slide.notes_text_frame
    notes.text = notes_text

# -------------------- SLIDE 1: TITLE --------------------
s = prs.slides.add_slide(BLANK); set_bg(s)
add_text(s, "Reading 23,207 Companies",
         Inches(0.8), Inches(2.3), Inches(11.7), Inches(1.4),
         size=54, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_rule(s, Inches(5.67), Inches(3.85), Inches(2.0))
add_text(s, "A capstone in machine intelligence, intellectual honesty, and library engineering.",
         Inches(1.0), Inches(4.1), Inches(11.3), Inches(0.8),
         size=22, color=LIME, italic=True, align=PP_ALIGN.CENTER)
add_text(s, "MGT 599  ·  Group 4  ·  Akash Anipakalu Giridhar  ·  May 18, 2026",
         Inches(1.0), Inches(6.5), Inches(11.3), Inches(0.5),
         size=14, color=DIM, align=PP_ALIGN.CENTER)
add_notes(s, "Morningstar has 23,000 public companies on its books, and every one needs to land in one of 145 industry boxes. Today we'll show you how we taught a machine to do that job — and what we found when we audited our own answer.")

# -------------------- SLIDE 2: CHAPTER 1 BEAT --------------------
s = prs.slides.add_slide(BLANK); set_bg(s)
add_text(s, "145", Inches(0.5), Inches(1.8), Inches(12.3), Inches(3.0),
         size=220, color=LIME, bold=True, align=PP_ALIGN.CENTER)
add_text(s, "industry codes. one company. one shot.",
         Inches(0.5), Inches(5.0), Inches(12.3), Inches(0.7),
         size=28, color=WHITE, align=PP_ALIGN.CENTER, italic=True)
add_footer(s, "Chapter 1 — The Problem")
add_notes(s, "Section header. Pause. Let the number breathe. ~15 seconds.")

# -------------------- SLIDE 3: THE PROBLEM --------------------
s = prs.slides.add_slide(BLANK); set_bg(s)
add_text(s, "The Classification Bottleneck",
         Inches(0.7), Inches(0.6), Inches(12.0), Inches(1.0),
         size=40, color=WHITE, bold=True)
add_rule(s, Inches(0.7), Inches(1.55), Inches(1.5))
add_bullets(s, [
    "Morningstar's GECS — 145 industries, 428 sub-industries",
    "23,207 companies (53,585 segment records) hand-classified by analysts",
    "Misclassification distorts sector ETFs, peer benchmarks, factor models",
    "A conglomerate may legitimately belong to 4+ codes at once",
], Inches(0.9), Inches(2.2), Inches(11.5), Inches(4.5), size=24)
add_footer(s, "Slide 3 — The Problem")
add_notes(s, "This isn't theoretical. Misclassification has real money attached — fund indices, risk models, peer comparisons all roll up by industry code. And the hardest cases are exactly the most valuable companies: the multi-segment conglomerates that don't fit one box.")

# -------------------- SLIDE 4: APPROACH + ROLES --------------------
s = prs.slides.add_slide(BLANK); set_bg(s)
add_text(s, "How Group 4 Attacked the Problem",
         Inches(0.7), Inches(0.5), Inches(12.0), Inches(1.0),
         size=36, color=WHITE, bold=True)
add_text(s, "Problem → Data → Library → Models → Audit → Path → Close",
         Inches(0.7), Inches(1.4), Inches(12.0), Inches(0.5),
         size=14, color=LIME, italic=True)
add_rule(s, Inches(0.7), Inches(1.9), Inches(12.0), thickness_pt=1)

# Left column — Technical
add_text(s, "Technical approach",
         Inches(0.9), Inches(2.2), Inches(5.7), Inches(0.5),
         size=18, color=LIME, bold=True)
add_bullets(s, [
    "TF-IDF baseline → ModernBERT-large → cascade",
    "Audit-first: company-disjoint splits",
    "breezeml library shipped to PyPI",
    "Flask + Next.js production stack",
], Inches(0.9), Inches(2.8), Inches(5.7), Inches(3.5), size=18)

# Right column — Team
add_text(s, "Team roles",
         Inches(7.0), Inches(2.2), Inches(5.7), Inches(0.5),
         size=18, color=LIME, bold=True)
add_bullets(s, [
    "Akash A.G. — architecture, breezeml, ModernBERT",
    "[Teammate 2] — data pipeline, leakage audit",
    "[Teammate 3] — frontend, demo infrastructure",
    "[Teammate 4] — evaluation, documentation",
], Inches(7.0), Inches(2.8), Inches(5.7), Inches(3.5), size=18)
add_footer(s, "Slide 4 — Approach + Roles")
add_notes(s, "Three principles guided us: build something we'd actually use, audit our own numbers harder than anyone else would, ship real code — not just a notebook. That brings us to the data.")

# -------------------- SLIDE 5: THE DATA --------------------
s = prs.slides.add_slide(BLANK); set_bg(s)
add_text(s, "What the Model Reads",
         Inches(0.7), Inches(0.6), Inches(12.0), Inches(1.0),
         size=40, color=WHITE, bold=True)
add_rule(s, Inches(0.7), Inches(1.55), Inches(1.5))
add_bullets(s, [
    "Source: task1_clean.csv — 23,207 companies / 53,585 segment records",
    "Input: LongProfile (free-form business description)",
    "Labels: 145 industry codes (Task 1), 428 sub-industry codes (Task 2)",
    "Enrichment: SegmentName, SegmentDescription, Revenue, revenue_share",
    "35.1% of companies are multi-segment → 55.2% of rows ambiguous by construction",
], Inches(0.9), Inches(2.2), Inches(11.5), Inches(4.5), size=22)
add_footer(s, "Slide 5 — The Data")
add_notes(s, "Over a third of companies don't have one answer — they have several, weighted by which segment makes the money. The model isn't just classifying; it's adjudicating. Which forced us to rebuild our tooling, starting with the library itself.")

# -------------------- SLIDE 6: BREEZEML ★ --------------------
s = prs.slides.add_slide(BLANK); set_bg(s)
add_text(s, "breezeml",
         Inches(0.7), Inches(0.5), Inches(12.0), Inches(1.0),
         size=54, color=WHITE, bold=True, font="Consolas")
add_text(s, "We didn't use a library. We shipped one.",
         Inches(0.7), Inches(1.45), Inches(12.0), Inches(0.6),
         size=22, color=LIME, italic=True)
# pip install pill
pill = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.5), Inches(2.4), Inches(4.3), Inches(0.7))
pill.fill.solid(); pill.fill.fore_color.rgb = BG
pill.line.color.rgb = LIME; pill.line.width = Pt(1.5)
tf = pill.text_frame; tf.margin_left = Inches(0.2); tf.margin_right = Inches(0.2)
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "pip install breezeml"
r.font.name = "Consolas"; r.font.size = Pt(22); r.font.color.rgb = LIME; r.font.bold = True

add_bullets(s, [
    "Zero-boilerplate ML framework built on scikit-learn — authored by Akash A.G.",
    "v0.2.3 — Primal SVM fix: training time 20+ min → under 2 sec",
    "v0.2.5 — Balanced class weights: eliminated SMOTE entirely",
    "Level 2 — 3-level cascade (Sector → Industry Group → Code)",
    "Five public PyPI releases shipped during this capstone",
], Inches(0.9), Inches(3.4), Inches(11.5), Inches(3.5), size=20)
add_footer(s, "Slide 6 — breezeml")
add_notes(s, "Most capstone teams use libraries. We shipped one. When LinearSVC deadlocked on our text data, we didn't switch models — we found that the dual SVM formulation was mathematically wrong for our data shape and patched the upstream library. That patch is live on PyPI now. Anyone in this room can install it. Which raises the question — how well did the models that used it actually perform?")

# -------------------- SLIDE 7: CHAPTER 2 BEAT --------------------
s = prs.slides.add_slide(BLANK); set_bg(s)
add_text(s, "70.29%", Inches(0.5), Inches(1.7), Inches(12.3), Inches(3.0),
         size=180, color=LIME, bold=True, align=PP_ALIGN.CENTER)
add_text(s, "what honest classification looks like.",
         Inches(0.5), Inches(4.9), Inches(12.3), Inches(0.7),
         size=28, color=WHITE, align=PP_ALIGN.CENTER, italic=True)
add_footer(s, "Chapter 2 — The Models")
add_notes(s, "Pause beat. ~10 seconds.")

# -------------------- SLIDE 8: F1 PROGRESSION CHART --------------------
s = prs.slides.add_slide(BLANK); set_bg(s)
add_text(s, "Four Models, Four Steps Up the Wall",
         Inches(0.7), Inches(0.5), Inches(12.0), Inches(1.0),
         size=36, color=WHITE, bold=True)
add_rule(s, Inches(0.7), Inches(1.45), Inches(1.5))
# Chart placeholder
ph = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(1.9), Inches(11.3), Inches(4.5))
ph.fill.solid(); ph.fill.fore_color.rgb = RGBColor(0x05, 0x1C, 0x1B)
ph.line.color.rgb = LIME; ph.line.width = Pt(0.75)
tf = ph.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "[ GRAPH 1 — F1 Progression bar chart ]\nlime + white-40% bars on #001514"
r.font.name = "Calibri"; r.font.size = Pt(18); r.font.color.rgb = DIM; r.font.italic = True

add_text(s, "Each jump smaller than the last. Something else was going on.",
         Inches(1.0), Inches(6.6), Inches(11.3), Inches(0.5),
         size=16, color=LIME, italic=True, align=PP_ALIGN.CENTER)
add_footer(s, "Slide 8 — Four Models")
add_notes(s, "We climbed from 59% to 70% by doing the right things — bigger models, better features, more careful training. But the jumps got smaller every time. By the time we hit ModernBERT-large, we were spending three hours of GPU to move the needle a single point. Something else was going on. And then we found it.")

# -------------------- SLIDE 9: CHAPTER 3 BEAT (CORAL) --------------------
s = prs.slides.add_slide(BLANK); set_bg(s)
add_text(s, "88.9%", Inches(0.5), Inches(1.7), Inches(12.3), Inches(3.0),
         size=200, color=CORAL, bold=True, align=PP_ALIGN.CENTER)
# Strikethrough simulation via line over text
line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.5), Inches(5.32), Inches(6.3), Pt(3))
line.fill.solid(); line.fill.fore_color.rgb = WHITE
line.line.fill.background()
add_text(s, "was not real.",
         Inches(0.5), Inches(5.0), Inches(12.3), Inches(0.8),
         size=36, color=WHITE, align=PP_ALIGN.CENTER)
add_footer(s, "Chapter 3 — The Audit")
add_notes(s, "Pause. Let the strikethrough land. ~15 seconds.")

# -------------------- SLIDE 10: LEAKAGE AUDIT --------------------
s = prs.slides.add_slide(BLANK); set_bg(s)
add_text(s, "The 88.9% That Wasn't Real",
         Inches(0.7), Inches(0.5), Inches(12.0), Inches(1.0),
         size=36, color=WHITE, bold=True)
add_rule(s, Inches(0.7), Inches(1.45), Inches(1.5))
# Donut placeholder (left)
ph = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), Inches(2.0), Inches(4.5), Inches(4.5))
ph.fill.solid(); ph.fill.fore_color.rgb = CORAL
ph.line.fill.background()
inner = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.8), Inches(3.0), Inches(2.5), Inches(2.5))
inner.fill.solid(); inner.fill.fore_color.rgb = BG
inner.line.fill.background()
add_text(s, "97.2%",
         Inches(1.8), Inches(3.5), Inches(2.5), Inches(1.0),
         size=36, color=CORAL, bold=True, align=PP_ALIGN.CENTER)
# Right column text
add_text(s, "We caught it.\nWe documented it.\nWe rebuilt the splits.",
         Inches(6.2), Inches(2.3), Inches(6.5), Inches(2.0),
         size=24, color=WHITE)
add_text(s, "CASCADE_AUDIT.md — the receipts",
         Inches(6.2), Inches(4.5), Inches(6.5), Inches(0.6),
         size=18, color=LIME, italic=True)
add_text(s, "Honest baseline since: 70.29% on company-disjoint test",
         Inches(6.2), Inches(5.2), Inches(6.5), Inches(0.6),
         size=18, color=WHITE)
add_footer(s, "Slide 10 — The Audit")
add_notes(s, "Early in the project our cascade reported 88.9%. We didn't trust it. We audited the splits and found that 97.2% of test rows had been seen by the model during training — the same company's text appearing on both sides of the split. We could have shipped the 88.9%. Nobody outside the team would have known. We didn't. We documented the leakage, rebuilt the data pipeline to be company-disjoint, and reported a lower honest number. That's not a setback. That's the most professional thing we did in this project.")

# -------------------- SLIDE 11: THE CEILING --------------------
s = prs.slides.add_slide(BLANK); set_bg(s)
add_text(s, "Why 80% Is a Wall",
         Inches(0.7), Inches(0.5), Inches(12.0), Inches(1.0),
         size=40, color=WHITE, bold=True)
add_rule(s, Inches(0.7), Inches(1.45), Inches(1.5))
# Stacked bar simulation
def seg(left_in, width_in, color, label, label_color=BG):
    rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left_in), Inches(2.5), Inches(width_in), Inches(1.2))
    rect.fill.solid(); rect.fill.fore_color.rgb = color
    rect.line.fill.background()
    tf = rect.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label
    r.font.name = "Calibri"; r.font.size = Pt(16); r.font.color.rgb = label_color; r.font.bold = True

# Widths proportional to 45.3, 25.5, 15.9, 13.3 — total ~100, scale to 11.5 in
total_w = 11.5
left = 0.9
for pct, color, lab, lc in [
    (45.3, LIME, "45.3% Single-code", BG),
    (25.5, RGBColor(0x91, 0x9C, 0x58), "25.5% 2-code", BG),
    (15.9, RGBColor(0x61, 0x68, 0x3B), "15.9% 3-code", WHITE),
    (13.3, CORAL, "13.3% 4+ code", BG),
]:
    w = total_w * pct / 100
    seg(left, w, color, lab, lc)
    left += w

add_text(s, "Even a perfect single-code classifier + 60% multi-code accuracy ≈ 76% macro F1",
         Inches(0.9), Inches(4.3), Inches(11.5), Inches(0.6),
         size=20, color=WHITE, italic=True, align=PP_ALIGN.CENTER)
add_text(s, "The ceiling is the data, not the model.",
         Inches(0.9), Inches(5.5), Inches(11.5), Inches(0.6),
         size=24, color=LIME, bold=True, align=PP_ALIGN.CENTER)
add_footer(s, "Slide 11 — The Ceiling")
add_notes(s, "Seventy percent honest is better than eighty-nine percent fake. And the math tells us why eighty is hard: a third of the test set is companies that legitimately belong to multiple codes. Even a perfect classifier can't be right about a company that has four right answers. So what do we do about it?")

# -------------------- SLIDE 12: CHAPTER 4 BEAT --------------------
s = prs.slides.add_slide(BLANK); set_bg(s)
add_text(s, "80%", Inches(0.5), Inches(1.7), Inches(12.3), Inches(3.0),
         size=220, color=LIME, bold=True, align=PP_ALIGN.CENTER)
add_text(s, "by term end. Here's how.",
         Inches(0.5), Inches(5.0), Inches(12.3), Inches(0.7),
         size=28, color=WHITE, italic=True, align=PP_ALIGN.CENTER)
add_footer(s, "Chapter 4 — What's Next")
add_notes(s, "Pause beat. ~10 seconds.")

# -------------------- SLIDE 13: FOUR PATHS --------------------
s = prs.slides.add_slide(BLANK); set_bg(s)
add_text(s, "Four Paths to 80%",
         Inches(0.7), Inches(0.5), Inches(12.0), Inches(1.0),
         size=40, color=WHITE, bold=True)
add_rule(s, Inches(0.7), Inches(1.45), Inches(1.5))

# Options A/B/D — normal
add_bullets(s, [
    "A.  Decidable-subset F1 — score only on unambiguous rows (defensible 80%+)",
    "B.  Revenue-weighted per-company prediction",
], Inches(0.9), Inches(2.0), Inches(11.5), Inches(1.6), size=20)

# Option C — highlighted box
boxC = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(3.5), Inches(11.9), Inches(1.4))
boxC.fill.solid(); boxC.fill.fore_color.rgb = BG
boxC.line.color.rgb = LIME; boxC.line.width = Pt(2)
tf = boxC.text_frame; tf.margin_left = Inches(0.3); tf.margin_top = Inches(0.15)
p = tf.paragraphs[0]
r = p.add_run(); r.text = "C.  Sector-conditioned head on ModernBERT-large embeddings"
r.font.name = "Calibri"; r.font.size = Pt(22); r.font.color.rgb = WHITE; r.font.bold = True
p2 = tf.add_paragraph()
r2 = p2.add_run(); r2.text = "Best legitimate path — projected 75–78%"
r2.font.name = "Calibri"; r2.font.size = Pt(18); r2.font.color.rgb = LIME; r2.font.italic = True

add_bullets(s, [
    "D.  Brute-force longer training — capped ~71–73%",
    "In flight: segment-aware + revenue-weighted + distilled + seed-123 ensemble variants",
], Inches(0.9), Inches(5.2), Inches(11.5), Inches(1.6), size=20)
add_footer(s, "Slide 13 — Path Forward")
add_notes(s, "We're not done. Additional training runs are in flight as we speak. The most promising — Option C — combines the hierarchy of the cascade with what ModernBERT learned about language. By week's end, 75 to 78 percent. By term end, a defensible 80. To close, the full picture.")

# -------------------- SLIDE 14: WHERE WE LAND --------------------
s = prs.slides.add_slide(BLANK); set_bg(s)
add_text(s, "Where We Land",
         Inches(0.7), Inches(0.5), Inches(12.0), Inches(1.0),
         size=40, color=WHITE, bold=True)
add_rule(s, Inches(0.7), Inches(1.45), Inches(1.5))

col_w = 4.0
col_y = 2.2

def col(label, items, x_in):
    add_text(s, label, Inches(x_in), Inches(col_y), Inches(col_w), Inches(0.6),
             size=22, color=LIME, bold=True)
    add_bullets(s, items, Inches(x_in), Inches(col_y + 0.7), Inches(col_w), Inches(4.0), size=18)

col("Built", [
    "ModernBERT-large at 70.29% honest macro F1",
    "breezeml — 5 PyPI releases + Level 2",
    "Production backend + Next.js dashboard",
], 0.7)
col("Learned", [
    "Audit your own numbers first",
    "The data ceiling beats the model",
    "Hierarchy-first beats end-to-end",
], 4.85)
col("Next", [
    "Push to 75–78% via Option C",
    "Finalize Task 2 sub-industry results",
    "Open-source company-disjoint splits",
], 9.0)
add_footer(s, "Slide 14 — Close")
add_notes(s, "We taught a model to read a company. We got to 70 percent honest, we built a library that's public infrastructure now, and along the way we taught ourselves what intellectual honesty looks like in machine learning. Thank you. Happy to take questions.")

# -------------------- SLIDE 15: THANK YOU --------------------
s = prs.slides.add_slide(BLANK); set_bg(s)
add_text(s, "Thank you.",
         Inches(0.5), Inches(2.5), Inches(12.3), Inches(1.5),
         size=72, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_rule(s, Inches(5.67), Inches(4.1), Inches(2.0))
add_text(s, "pip install breezeml",
         Inches(0.5), Inches(4.3), Inches(12.3), Inches(0.8),
         size=26, color=LIME, align=PP_ALIGN.CENTER, font="Consolas")
add_text(s, "github.com/venomez-viper/Classification-Project",
         Inches(0.5), Inches(5.5), Inches(12.3), Inches(0.5),
         size=14, color=DIM, align=PP_ALIGN.CENTER)
add_text(s, "pypi.org/project/breezeml",
         Inches(0.5), Inches(5.9), Inches(12.3), Inches(0.5),
         size=14, color=DIM, align=PP_ALIGN.CENTER)
add_text(s, "Q&A →",
         Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.5),
         size=16, color=LIME, italic=True, align=PP_ALIGN.CENTER)
add_notes(s, "Open for questions.")

# -------------------- SAVE --------------------
out = r"C:\Users\akash\Desktop\capstone MGT 599\REFERENCE_DECK.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
